# -*- coding: utf-8 -*-
"""
Created on Fri Aug 29 12:07:20 2014

@author: ecksjoh
"""

from pptx import Presentation
from pptx.util import Inches, Px

img_path = 'monty-truth.png'

prs = Presentation()
blank_slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_slide_layout)

left = top = Inches(1)
pic = slide.shapes.add_picture(img_path, left, top)

left = Inches(5)
width = Px(280)
height = int(width*1.427)
pic = slide.shapes.add_picture(img_path, left, top, width, height)

prs.save('test.pptx')
